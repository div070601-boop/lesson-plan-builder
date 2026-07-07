"""
PPTX Builder Service — Assembles generated content into PowerPoint files.
Uses python-pptx with slide master template support (PRD Section 7, Step 5).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from typing import Optional

from models.generation import SlideContent, SlideType


class PPTXBuilder:
    """Builds PPTX files from generated slide content."""

    def __init__(self, template_path: Optional[str] = None):
        """Initialize with an optional slide master template."""
        if template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            # Set widescreen 16:9 aspect ratio
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)

    def build(self, slides: list[SlideContent], output_path: str) -> str:
        """Build the full PPTX from slide content list."""
        for slide_content in slides:
            self._add_slide(slide_content)

        # Ensure output directory exists
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        return output_path

    def _add_slide(self, content: SlideContent) -> None:
        """Add a single slide based on its type and content."""
        if content.slide_type == SlideType.TITLE:
            self._add_title_slide(content)
        elif content.slide_type in (SlideType.OBJECTIVES, SlideType.AGENDA):
            self._add_bullet_slide(content)
        elif content.slide_type == SlideType.ACTIVITY:
            self._add_activity_slide(content)
        elif content.slide_type == SlideType.SUMMARY:
            self._add_bullet_slide(content)
        elif content.slide_type == SlideType.QUIZ:
            self._add_bullet_slide(content)
        else:
            self._add_content_slide(content)

    def _add_title_slide(self, content: SlideContent) -> None:
        """Add a title slide."""
        layout = self.prs.slide_layouts[0]  # Title Slide layout
        slide = self.prs.slides.add_slide(layout)

        if slide.placeholders:
            if 0 in slide.placeholders:
                slide.placeholders[0].text = content.title
            if 1 in slide.placeholders and content.body:
                slide.placeholders[1].text = "\n".join(content.body)

        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _add_content_slide(self, content: SlideContent) -> None:
        """Add a standard content slide with title and bullets."""
        layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(layout)

        if 0 in slide.placeholders:
            slide.placeholders[0].text = content.title

        if 1 in slide.placeholders and content.body:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, bullet in enumerate(content.body):
                if i == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _add_bullet_slide(self, content: SlideContent) -> None:
        """Add a slide with bullet points (objectives, agenda, summary)."""
        self._add_content_slide(content)  # Same layout, different styling later

    def _add_activity_slide(self, content: SlideContent) -> None:
        """Add an activity slide with instructions and timing."""
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)

        if 0 in slide.placeholders:
            slide.placeholders[0].text = f"[Activity] {content.title}"

        if 1 in slide.placeholders and content.body:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, bullet in enumerate(content.body):
                if i == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

            if content.estimated_duration:
                p = tf.add_paragraph()
                p.text = f"\n⏱ Duration: {content.estimated_duration}"
                p.level = 0

        # Combine speaker notes with activity instructions
        notes = []
        if content.speaker_notes:
            notes.append(content.speaker_notes)
        if content.activity_instructions:
            notes.append(f"\nActivity Setup:\n{content.activity_instructions}")

        if notes:
            slide.notes_slide.notes_text_frame.text = "\n".join(notes)


pptx_builder = PPTXBuilder()
